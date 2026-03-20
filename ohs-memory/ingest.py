"""
OHS Memory — Document Ingestion Script
=======================================
Ingests PDF, Word, PowerPoint, Excel, Markdown, plain text, and Outlook .msg
(email) files into the OHS organizational memory database.

Supported formats:
  .pdf   — text extracted via Claude API (handles scans, slide decks, etc.)
  .docx  — Word documents via python-docx (paragraphs + tables)
  .pptx  — PowerPoint via python-pptx (all slides + speaker notes)
  .xlsx  — Excel via openpyxl (all sheets, row-by-row)
  .md / .txt — read directly
  .msg   — Outlook email via extract-msg (headers + body + inline attachments)

Usage:
  python ingest.py path/to/document.pdf
  python ingest.py path/to/folder/           # batch-ingest entire folder (recursive)
  python ingest.py notes.pdf --subject "All-Staff" --year 2025-26 --type meeting_notes
  python ingest.py emails/ --subject "All-Staff" --type email --unit "Orion Planning Team"

Batch flags:
  --limit N      Only process the first N files (test before full run)
  --log FILE     Append failure details to FILE (default: ingest_errors.log)

For .msg files: --teacher and --year are auto-extracted from each email's
headers (sender display name and sent date). You can override with explicit flags.
"""

import argparse
import base64
import hashlib
import os
import re
import sys
import tempfile
from datetime import datetime, timezone
from email.utils import parsedate_to_datetime
from pathlib import Path

import psycopg2
import tiktoken
from dotenv import load_dotenv
from openai import OpenAI
from tqdm import tqdm

load_dotenv()

# ── Configuration ──────────────────────────────────────────────────────────────

SMALL_CHUNK_TOKENS  = 200   # precise, factual retrieval
SMALL_CHUNK_OVERLAP = 40
LARGE_CHUNK_TOKENS  = 900   # contextual retrieval
LARGE_CHUNK_OVERLAP = 120

# Claude model for PDF text extraction
EXTRACTION_MODEL = "claude-haiku-4-5"

# OpenAI embedding model — MUST match mcp_server.py and search-proxy.php
EMBEDDING_MODEL  = "text-embedding-3-small"
EMBEDDING_DIM    = 1536

# Orion HS opened: first day of school, fall 2024
# Emails before this date are tagged "pre-opening"
OHS_OPENING_DATE = datetime(2024, 9, 1, tzinfo=timezone.utc)

# Attachments larger than this (in tokens) get a truncated inline preview
# instead of full inline text, with a note to ingest separately.
ATTACHMENT_INLINE_TOKEN_LIMIT = 3000

# Supported file extensions
SUPPORTED = {".pdf", ".docx", ".pptx", ".xlsx", ".md", ".txt", ".msg"}

# ── Client initialization ──────────────────────────────────────────────────────

_anthropic = None
_openai: OpenAI | None = None
_tokenizer = None


def get_anthropic():
    global _anthropic
    if _anthropic is None:
        import anthropic as _lib
        _anthropic = _lib.Anthropic(api_key=os.environ["ANTHROPIC_API_KEY"])
    return _anthropic


def get_openai():
    global _openai
    if _openai is None:
        _openai = OpenAI(api_key=os.environ["OPENAI_API_KEY"])
    return _openai


def get_tokenizer():
    global _tokenizer
    if _tokenizer is None:
        _tokenizer = tiktoken.get_encoding("cl100k_base")
    return _tokenizer


# ── Date parsing ───────────────────────────────────────────────────────────────

def _parse_msg_date(raw) -> datetime | None:
    """
    Normalize whatever extract-msg returns for msg.date into a datetime (UTC).

    extract-msg may return:
      - a datetime object  (modern versions, well-formed dates)
      - a string           (unparseable RFC 2822, or older library versions)
      - None               (date header missing entirely)
    """
    if raw is None:
        return None
    if isinstance(raw, datetime):
        if raw.tzinfo is None:
            return raw.replace(tzinfo=timezone.utc)
        return raw
    # It's a string — try RFC 2822 first (the standard email date format)
    try:
        return parsedate_to_datetime(str(raw))
    except Exception:
        pass
    # Fallback: try a handful of common strptime patterns
    for fmt in (
        "%Y-%m-%d %H:%M:%S%z",
        "%Y-%m-%d %H:%M:%S",
        "%a, %d %b %Y %H:%M:%S %z",
        "%a, %d %b %Y %H:%M:%S",
    ):
        try:
            dt = datetime.strptime(str(raw).strip(), fmt)
            if dt.tzinfo is None:
                dt = dt.replace(tzinfo=timezone.utc)
            return dt
        except ValueError:
            pass
    # Couldn't parse — treat as unknown
    return None


# ── School-year detection ──────────────────────────────────────────────────────

def detect_school_year(dt: datetime | None) -> str | None:
    """
    Map an email's sent date to the appropriate school-year tag.

    pre-opening  — before Sep 2024 (planning phase, school not yet open)
    2024-25      — Sep 2024 through Jun 2025 (first year)
    2025-26      — Jul 2025 onwards (second year)
    """
    if dt is None:
        return None
    if dt.tzinfo is None:
        dt = dt.replace(tzinfo=timezone.utc)
    if dt < OHS_OPENING_DATE:
        return "pre-opening"
    y, m = dt.year, dt.month
    if (y == 2024 and m >= 9) or (y == 2025 and m <= 6):
        return "2024-25"
    if (y == 2025 and m >= 7) or y == 2026:
        return "2025-26"
    # Generic fallback for future years
    return f"{y}-{str(y + 1)[-2:]}" if m >= 7 else f"{y - 1}-{str(y)[-2:]}"


# ── Text extraction — one function per format ─────────────────────────────────

def extract_text(path: Path) -> tuple[str, dict]:
    """
    Dispatch to the appropriate extractor based on file extension.
    Returns (text, extra_metadata).

    extra_metadata is populated from file headers where available (.msg only),
    providing auto-detected teacher and school_year fields.
    """
    suffix = path.suffix.lower()
    if suffix in (".md", ".txt"):
        return path.read_text(encoding="utf-8", errors="replace"), {}
    if suffix == ".pdf":
        return _extract_text_from_pdf(path), {}
    if suffix == ".docx":
        return _extract_text_from_docx(path), {}
    if suffix == ".pptx":
        return _extract_text_from_pptx(path), {}
    if suffix == ".xlsx":
        return _extract_text_from_xlsx(path), {}
    if suffix == ".msg":
        return _extract_text_from_msg(path)
    raise ValueError(f"Unsupported file type: {suffix}")


def _extract_text_from_pdf(path: Path) -> str:
    """Send PDF to Claude API for full text extraction."""
    with open(path, "rb") as f:
        pdf_b64 = base64.standard_b64encode(f.read()).decode("utf-8")

    print(f"  Extracting PDF via Claude API ({EXTRACTION_MODEL})...")

    with get_anthropic().messages.stream(
        model=EXTRACTION_MODEL,
        max_tokens=8192,
        messages=[{
            "role": "user",
            "content": [
                {
                    "type": "document",
                    "source": {"type": "base64", "media_type": "application/pdf", "data": pdf_b64},
                },
                {
                    "type": "text",
                    "text": (
                        "Extract ALL text from this document. "
                        "Preserve structure: headings, bullet points, numbered lists, tables. "
                        "If this is a slide deck, include ALL slide content. "
                        "Do NOT summarize — extract the complete text verbatim. "
                        "Return ONLY the extracted text with no commentary."
                    ),
                },
            ],
        }],
    ) as stream:
        final = stream.get_final_message()

    return final.content[0].text


def _extract_text_from_docx(path: Path) -> str:
    """
    Extract text from a Word document.
    Captures paragraphs (with heading labels) and table cell content.
    Requires: pip install python-docx
    """
    try:
        from docx import Document
    except ImportError:
        raise ImportError("python-docx required for .docx files.  Run: pip install python-docx")

    print("  Extracting DOCX...")
    doc = Document(str(path))
    parts = []

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        # Prefix headings so structure is visible in chunks
        if para.style.name.startswith("Heading"):
            level = para.style.name.split()[-1] if para.style.name.split()[-1].isdigit() else "1"
            parts.append(f"{'#' * int(level)} {text}")
        else:
            parts.append(text)

    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                parts.append(" | ".join(cells))

    return "\n".join(parts)


def _extract_text_from_pptx(path: Path) -> str:
    """
    Extract text from a PowerPoint presentation.
    Captures all text shapes and speaker notes per slide.
    Requires: pip install python-pptx
    """
    try:
        from pptx import Presentation
    except ImportError:
        raise ImportError("python-pptx required for .pptx files.  Run: pip install python-pptx")

    print("  Extracting PPTX...")
    prs = Presentation(str(path))
    parts = []

    for i, slide in enumerate(prs.slides, 1):
        slide_texts = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                t = "".join(run.text for run in para.runs).strip()
                if t:
                    slide_texts.append(t)

        if slide_texts:
            parts.append(f"\n--- Slide {i} ---")
            parts.extend(slide_texts)

        # Speaker notes are often the richest text in a teacher's deck
        try:
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    parts.append(f"[Speaker notes: {notes}]")
        except Exception:
            pass

    return "\n".join(parts)


def _extract_text_from_xlsx(path: Path) -> str:
    """
    Extract text from an Excel workbook.
    Outputs each sheet as a tab-separated table; empty rows are skipped.
    Requires: pip install openpyxl
    """
    try:
        import openpyxl
    except ImportError:
        raise ImportError("openpyxl required for .xlsx files.  Run: pip install openpyxl")

    print("  Extracting XLSX...")
    wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    parts = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        parts.append(f"\n=== Sheet: {sheet_name} ===")
        row_count = 0
        for row in ws.iter_rows(values_only=True):
            cells = [str(v).strip() for v in row if v is not None and str(v).strip()]
            if cells:
                parts.append("\t".join(cells))
                row_count += 1
        if row_count == 0:
            parts.append("(empty sheet)")

    wb.close()
    return "\n".join(parts)


def _extract_text_from_msg(path: Path) -> tuple[str, dict]:
    """
    Extract text and metadata from an Outlook .msg file.

    Returns (text, extra_metadata) where extra_metadata may contain:
      teacher     — sender display name
      school_year — auto-detected from sent date
      notes       — attachment manifest

    PDF and DOCX attachments are extracted inline (up to ATTACHMENT_INLINE_TOKEN_LIMIT).
    Larger attachments get a truncated preview; unsupported types are noted.

    Requires: pip install extract-msg
    """
    try:
        import extract_msg as em
    except ImportError:
        raise ImportError("extract-msg required for .msg files.  Run: pip install extract-msg")

    with em.Message(str(path)) as msg:
        subject  = (msg.subject or "(no subject)").strip()
        sender   = (msg.sender  or "").strip()
        to_field = (msg.to      or "").strip()
        sent_dt  = _parse_msg_date(msg.date)

        body = msg.body or ""
        if not body.strip() and msg.htmlBody:
            body = re.sub(r"<[^>]+>", " ", msg.htmlBody or "")
            body = re.sub(r"\s{2,}", " ", body).strip()

        parts = []
        if subject:  parts.append(f"Subject: {subject}")
        if sender:   parts.append(f"From: {sender}")
        if to_field: parts.append(f"To: {to_field}")
        if sent_dt:  parts.append(f"Date: {sent_dt.strftime('%Y-%m-%d %H:%M')}")
        parts.append("")
        parts.append(body.strip())

        # ── Inline attachment extraction ──────────────────────────────────────
        attachment_notes = []

        for att in (msg.attachments or []):
            fname = (att.longFilename or att.shortFilename or "").strip()
            if not fname:
                continue

            suffix = Path(fname).suffix.lower()
            att_data = att.data

            if not att_data:
                attachment_notes.append(f"{fname} (empty, skipped)")
                continue

            # Extract text for supported types
            att_text = None
            if suffix == ".pdf":
                try:
                    att_text = _extract_attachment_bytes(att_data, suffix, fname)
                except Exception as e:
                    attachment_notes.append(f"{fname} (PDF extraction failed: {e})")
            elif suffix == ".docx":
                try:
                    att_text = _extract_attachment_bytes(att_data, suffix, fname)
                except Exception as e:
                    attachment_notes.append(f"{fname} (DOCX extraction failed: {e})")
            elif suffix in (".txt", ".md"):
                att_text = att_data.decode("utf-8", errors="replace")
            else:
                # Images, XLSX, PPTX, etc. — note existence, don't extract
                label = suffix[1:].upper() if suffix else "file"
                attachment_notes.append(
                    f"{fname} ({label} — not extracted inline; "
                    f"save from email and ingest separately if needed)"
                )
                parts.append(f"\n=== Attachment: {fname} ===")
                parts.append(f"[{label} attached to this email. "
                              f"Content not extracted. Source: email only.]")
                continue

            if att_text is None:
                continue

            # Inline if small enough; truncate with note if large
            enc = get_tokenizer()
            tokens = enc.encode(att_text)

            if len(tokens) <= ATTACHMENT_INLINE_TOKEN_LIMIT:
                parts.append(f"\n=== Attachment: {fname} ({len(tokens)} tokens) ===")
                parts.append(att_text)
                attachment_notes.append(f"{fname} (fully extracted, {len(tokens)} tokens)")
            else:
                preview = enc.decode(tokens[:ATTACHMENT_INLINE_TOKEN_LIMIT])
                parts.append(f"\n=== Attachment: {fname} (preview — {len(tokens)} tokens total) ===")
                parts.append(preview)
                parts.append(f"\n[... {len(tokens) - ATTACHMENT_INLINE_TOKEN_LIMIT} tokens truncated. "
                              f"Ingest {fname} separately for full text.]")
                attachment_notes.append(
                    f"{fname} ({len(tokens)} tokens — preview only, ingest separately)"
                )

        text = "\n".join(parts)

        # ── Auto-extract metadata from headers ────────────────────────────────
        extra: dict = {}
        sender_name = re.sub(r"<[^>]+>", "", sender).strip().strip('"')
        if sender_name:
            extra["teacher"] = sender_name
        if sent_dt:
            extra["school_year"] = detect_school_year(sent_dt)
        if attachment_notes:
            extra["notes"] = "Attachments: " + "; ".join(attachment_notes)

    return text, extra


def _extract_attachment_bytes(data: bytes, suffix: str, fname: str) -> str:
    """Write attachment bytes to a temp file and extract text from it."""
    print(f"    Extracting attachment: {fname}...")
    tmp_path = None
    try:
        with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as tmp:
            tmp.write(data)
            tmp_path = Path(tmp.name)
        text, _ = extract_text(tmp_path)
        return text
    finally:
        if tmp_path and tmp_path.exists():
            tmp_path.unlink()


# ── Chunking & embedding ──────────────────────────────────────────────────────

def chunk_text(text: str, max_tokens: int, overlap: int) -> list[str]:
    enc = get_tokenizer()
    tokens = enc.encode(text)
    if not tokens:
        return []
    chunks, start = [], 0
    while start < len(tokens):
        end = min(start + max_tokens, len(tokens))
        chunks.append(enc.decode(tokens[start:end]))
        start += max(max_tokens - overlap, 1)
    return chunks


def embed_texts(texts: list[str]) -> list[list[float]]:
    all_embeddings = []
    for i in range(0, len(texts), 100):
        batch = texts[i : i + 100]
        response = get_openai().embeddings.create(input=batch, model=EMBEDDING_MODEL)
        all_embeddings.extend([item.embedding for item in response.data])
    return all_embeddings


def format_vector(v: list[float]) -> str:
    return "[" + ",".join(f"{x:.8f}" for x in v) + "]"


# ── Interactive metadata prompt ───────────────────────────────────────────────

def prompt_metadata(filename: str) -> dict:
    print(f"\n  Metadata for: {filename}")
    print("  (Press Enter to skip any field)\n")

    def ask(prompt, options=None):
        if options:
            print(f"     Options: {', '.join(options)}")
        value = input(f"  {prompt}: ").strip()
        return value if value else None

    subject     = ask("Subject (e.g. Economics, Physics, All-Staff)")
    school_year = ask("School year (e.g. 2025-26, pre-opening)")
    teacher     = ask("Teacher/author (e.g. Weisenfeld)")
    doc_type    = ask("Document type", ["lesson_plan", "meeting_notes", "policy", "email", "other"])
    if doc_type and doc_type not in ("lesson_plan", "meeting_notes", "policy", "email", "other"):
        doc_type = "other"
    unit  = ask("Unit/topic (e.g. Supply and Demand)")
    notes = ask("Any notes about this document")
    return {
        "subject": subject, "school_year": school_year, "teacher": teacher,
        "doc_type": doc_type or "other", "unit": unit, "notes": notes,
    }


# ── Core ingestion pipeline ───────────────────────────────────────────────────

def hash_file(path: Path) -> str:
    h = hashlib.sha256()
    with open(path, "rb") as f:
        for block in iter(lambda: f.read(65536), b""):
            h.update(block)
    return h.hexdigest()


def ingest_file(path: Path, metadata: dict | None = None,
                conn=None, db_url: str | None = None) -> str | None:
    """
    Full ingestion pipeline for a single file.
    Returns document UUID on success, None if skipped (already ingested).

    Pass an open psycopg2 connection via `conn` for batch efficiency —
    the function commits per-document but does NOT close a shared connection.
    If conn is None, a new connection is opened and closed per call.
    """
    own_conn = conn is None
    if own_conn:
        conn = psycopg2.connect(db_url or os.environ["DATABASE_URL"])

    cur = conn.cursor()
    try:
        # 1. Hash check — skip if already ingested
        file_hash = hash_file(path)
        cur.execute("SELECT id FROM documents WHERE source_hash = %s", (file_hash,))
        if existing := cur.fetchone():
            print(f"  [skip] {path.name} (id: {existing[0]})")
            return None

        # 2. Extract text
        raw_text, auto_meta = extract_text(path)
        token_count = len(get_tokenizer().encode(raw_text))
        print(f"  Extracted {token_count:,} tokens")

        # 3. Resolve metadata: auto_meta fills gaps in CLI metadata
        if metadata is None:
            if path.suffix.lower() == ".msg" and auto_meta:
                resolved = {
                    "subject": None, "school_year": auto_meta.get("school_year"),
                    "teacher": auto_meta.get("teacher"), "doc_type": "email",
                    "unit": None, "notes": auto_meta.get("notes"),
                }
            else:
                resolved = prompt_metadata(path.name)
        else:
            resolved = dict(metadata)
            for key, val in auto_meta.items():
                if not resolved.get(key):
                    resolved[key] = val

        # 4. Insert document record
        cur.execute(
            """
            INSERT INTO documents
                (source_hash, source_file, original_filename,
                 subject, school_year, teacher, doc_type, unit, notes, raw_text)
            VALUES (%s,%s,%s,%s,%s,%s,%s,%s,%s,%s)
            RETURNING id
            """,
            (
                file_hash, str(path.resolve()), path.name,
                resolved.get("subject"), resolved.get("school_year"), resolved.get("teacher"),
                resolved.get("doc_type", "other"), resolved.get("unit"),
                resolved.get("notes"), raw_text,
            ),
        )
        doc_id = cur.fetchone()[0]

        # 5. Chunk → embed → store (small + large)
        enc = get_tokenizer()
        for label, max_tok, overlap in [
            ("small", SMALL_CHUNK_TOKENS,  SMALL_CHUNK_OVERLAP),
            ("large", LARGE_CHUNK_TOKENS,  LARGE_CHUNK_OVERLAP),
        ]:
            chunks = chunk_text(raw_text, max_tok, overlap)
            if not chunks:
                continue
            print(f"  Embedding {len(chunks)} {label} chunks...")
            embeddings = embed_texts(chunks)
            cur.executemany(
                """
                INSERT INTO chunks (document_id, chunk_size, content, embedding, position, token_count)
                VALUES (%s,%s,%s,%s::vector,%s,%s)
                """,
                [
                    (doc_id, label, chunk, format_vector(emb), i, len(enc.encode(chunk)))
                    for i, (chunk, emb) in enumerate(zip(chunks, embeddings))
                ],
            )
            print(f"  [ok] {len(chunks)} {label} chunks stored")

        conn.commit()
        yr = resolved.get("school_year") or "?"
        tc = resolved.get("teacher") or "?"
        print(f"\n  [DONE] {path.name}  [{yr} / {tc}]  (id: {doc_id})\n")
        return str(doc_id)

    except Exception as e:
        conn.rollback()
        print(f"\n  [FAIL] {path.name} — {e}\n")
        raise
    finally:
        cur.close()
        if own_conn:
            conn.close()


# ── CLI ────────────────────────────────────────────────────────────────────────

def main():
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")

    parser = argparse.ArgumentParser(
        description="Ingest documents into the OHS Memory database.",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Supported formats: .pdf  .docx  .pptx  .xlsx  .md  .txt  .msg

Examples:
  python ingest.py meeting_notes.pdf --subject "All-Staff" --year 2025-26 --type meeting_notes
  python ingest.py lesson.pptx       --subject Physics     --year 2025-26 --type lesson_plan

  # Test first 10 emails before committing to the full 1400
  python ingest.py emails/ --subject "All-Staff" --type email --unit "Orion Planning Team" --limit 10

  # Full email flood — safe to stop and restart, already-ingested files are skipped
  python ingest.py emails/ --subject "All-Staff" --type email --unit "Orion Planning Team"

  # SharePoint folder via OneDrive sync
  python ingest.py "C:\\Users\\johnw\\OneDrive - PSD1\\Orion HS SharePoint\\Shared Documents\\"
        """,
    )
    parser.add_argument("path",     help="File or folder to ingest")
    parser.add_argument("--subject",help="Subject area (e.g. All-Staff, Physics, Economics)")
    parser.add_argument("--year",   help="School year (e.g. 2025-26, pre-opening). "
                                         "Auto-detected per .msg file if omitted.")
    parser.add_argument("--teacher",help="Teacher/author. Auto-detected per .msg if omitted.")
    parser.add_argument("--type", dest="doc_type",
        choices=["lesson_plan", "meeting_notes", "policy", "email", "other"])
    parser.add_argument("--unit",   help="Unit, topic, or source group (e.g. 'Orion Planning Team')")
    parser.add_argument("--limit",  type=int, default=0,
                        help="Only process first N files (0 = no limit). Useful for test runs.")
    parser.add_argument("--log",    default="ingest_errors.log",
                        help="File to append failure details to (default: ingest_errors.log)")
    args = parser.parse_args()

    cli_metadata = None
    if any([args.subject, args.year, args.teacher, args.doc_type, args.unit]):
        cli_metadata = {
            "subject":     args.subject,
            "school_year": args.year,
            "teacher":     args.teacher,
            "doc_type":    args.doc_type or "other",
            "unit":        args.unit,
            "notes":       None,
        }

    target = Path(args.path)

    if target.is_file():
        if target.suffix.lower() not in SUPPORTED:
            print(f"Unsupported type: {target.suffix}  (supported: {', '.join(sorted(SUPPORTED))})")
            sys.exit(1)
        ingest_file(target, metadata=cli_metadata)
        return

    if not target.is_dir():
        print(f"Error: {target} is not a file or directory")
        sys.exit(1)

    # ── Batch mode ────────────────────────────────────────────────────────────
    files = sorted(f for f in target.glob("**/*") if f.suffix.lower() in SUPPORTED)
    if not files:
        exts = ", ".join(sorted(SUPPORTED))
        print(f"No supported files ({exts}) found in {target}")
        sys.exit(1)

    if args.limit:
        files = files[: args.limit]
        print(f"--limit {args.limit}: processing first {len(files)} file(s)\n")
    else:
        print(f"Found {len(files)} file(s) in {target}\n")

    # Single shared DB connection for the whole batch — much faster than
    # open/close per file for large floods (1 connection vs 1400)
    db_url = os.environ["DATABASE_URL"]
    conn = psycopg2.connect(db_url)
    print(f"DB connection open. Starting batch...\n")

    ok = skip = fail = 0
    log_path = Path(args.log)

    try:
        for f in tqdm(files, desc="Ingesting", unit="file"):
            try:
                # Auto-derive unit from top-level subfolder when --unit not given.
                # e.g. target=".../Orion Planning Team - Documents/"
                #      file  =".../Bell Schedules/somefile.docx"
                #      → unit = "Bell Schedules"
                file_meta = cli_metadata
                if file_meta is not None and not file_meta.get("unit"):
                    try:
                        rel_parts = f.relative_to(target).parts
                        if len(rel_parts) > 1:
                            file_meta = {**file_meta, "unit": rel_parts[0]}
                    except ValueError:
                        pass

                result = ingest_file(f, metadata=file_meta, conn=conn)
                if result:
                    ok += 1
                else:
                    skip += 1
            except Exception as e:
                fail += 1
                with open(log_path, "a", encoding="utf-8") as log:
                    log.write(f"{datetime.now().isoformat()}  FAIL  {f}  —  {e}\n")
    finally:
        conn.close()

    print(f"\n{'='*60}")
    print(f"Batch complete:")
    print(f"  Ingested : {ok}")
    print(f"  Skipped  : {skip}  (already in DB)")
    print(f"  Failed   : {fail}  (see {log_path})")
    print(f"{'='*60}")


if __name__ == "__main__":
    main()
